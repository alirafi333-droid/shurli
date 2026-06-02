from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BG    = RGBColor(0x07, 0x09, 0x13)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xD0)
RED        = RGBColor(0xEF, 0x44, 0x44)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)

def add_bg(slide, color=DARK_BG):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_rect(slide, x, y, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return txb

def add_slide(layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_rect(s, 0, 7.42, 13.33, 0.08, CYAN)
# Decorative card suits
add_text(s, "♠  ♥  ♦  ♣", 0, 0.5, 13.33, 1.2, size=48, color=RGBColor(0x1E,0x1E,0x3A), align=PP_ALIGN.CENTER)
add_text(s, "🎲 GRAVITY FLIP POKER", 0, 1.6, 13.33, 1.5, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "A Full-Stack Software Engineering Project", 0, 3.1, 13.33, 0.8, size=26, color=CYAN, align=PP_ALIGN.CENTER, italic=True)
add_rect(s, 3.5, 4.1, 6.33, 0.04, GOLD)
add_text(s, "Process Models  •  Agile  •  Refactoring  •  Testing  •  Deployment", 0, 4.3, 13.33, 0.6, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Software Construction & Development — Final Project", 0, 6.5, 13.33, 0.6, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: Project Overview ─────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "PROJECT OVERVIEW", 0.5, 0.18, 12, 0.7, size=32, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 3.0, 0.04, GOLD)

points = [
    ("🎯", "Objective", "Convert a legacy CLI poker game into a fully-featured, multi-platform software product."),
    ("🗂️", "Tech Stack", "Python · Flask · SQLite · Vanilla JS/CSS · PyQt5 · Streamlit"),
    ("🚀", "Deliverables", "Web App · Native Desktop App · Streamlit Dashboard · REST API · Unit Tests"),
    ("📚", "SE Concepts", "Agile · SPI · Refactoring · CI/CD · Exception Handling · Peer Reviews"),
]
for i, (icon, title, desc) in enumerate(points):
    row_y = 1.3 + i * 1.4
    add_rect(s, 0.5, row_y, 12.3, 1.2, RGBColor(0x10,0x14,0x28))
    add_text(s, icon, 0.7, row_y + 0.1, 0.8, 0.9, size=30, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.6, row_y + 0.05, 2.5, 0.5, size=18, bold=True, color=GOLD)
    add_text(s, desc, 1.6, row_y + 0.55, 10.8, 0.5, size=14, color=LIGHT_GRAY)

# ── SLIDE 3: Process Model – Agile ──────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "PROCESS MODEL: AGILE (SCRUM)", 0.5, 0.18, 12, 0.7, size=32, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 4.5, 0.04, GOLD)

sprints = [
    ("Sprint 1", "CLI → REST API", "Refactored legacy terminal code into Flask endpoints", PURPLE),
    ("Sprint 2", "Web Frontend", "HTML/CSS/JS with 3D cards, glassmorphism, audio synthesis", CYAN),
    ("Sprint 3", "Desktop App", "PyQt5 + QWebEngine native wrapper with server polling", GREEN),
    ("Sprint 4", "Python-Only UI", "Streamlit dashboard; unit tests & CI/CD GitHub Actions", GOLD),
]
for i, (sprint, title, desc, color) in enumerate(sprints):
    col = i % 2; row = i // 2
    x = 0.5 + col * 6.4; y = 1.3 + row * 2.8
    add_rect(s, x, y, 6.0, 2.5, RGBColor(0x10,0x14,0x28))
    add_rect(s, x, y, 6.0, 0.45, color)
    add_text(s, sprint, x + 0.15, y + 0.05, 5.5, 0.35, size=14, bold=True, color=DARK_BG)
    add_text(s, title, x + 0.15, y + 0.55, 5.7, 0.45, size=18, bold=True, color=WHITE)
    add_text(s, desc, x + 0.15, y + 1.1, 5.7, 1.2, size=13, color=LIGHT_GRAY)

# ── SLIDE 4: Architecture ────────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "SYSTEM ARCHITECTURE", 0.5, 0.18, 12, 0.7, size=32, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 4.0, 0.04, GOLD)

layers = [
    ("PRESENTATION LAYER", [("🌐 Web Browser", "index.html + app.js"), ("🖥️ Desktop GUI", "PyQt5 + QWebEngine"), ("📊 Streamlit", "streamlit_app.py")], PURPLE),
    ("API LAYER", [("🔗 REST Endpoints", "/api/place_bet, /api/fold"), ("🔒 Session Mgmt", "Flask Session + SQLite"), ("⚡ Flask Server", "Port 5000/5055")], CYAN),
    ("CORE ENGINE", [("🃏 HandEvaluator", "7-card best hand algorithm"), ("🤖 AIPlayer", "Probabilistic decision engine"), ("🗄️ Database", "SQLite ORM wrapper")], GREEN),
]
for i, (layer, items, color) in enumerate(layers):
    x = 0.4 + i * 4.3
    add_rect(s, x, 1.2, 4.0, 5.9, RGBColor(0x0C,0x10,0x22))
    add_rect(s, x, 1.2, 4.0, 0.5, color)
    add_text(s, layer, x + 0.1, 1.25, 3.8, 0.4, size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    for j, (name, detail) in enumerate(items):
        iy = 1.9 + j * 1.65
        add_rect(s, x + 0.15, iy, 3.7, 1.4, RGBColor(0x12,0x16,0x2E))
        add_text(s, name, x + 0.25, iy + 0.1, 3.5, 0.45, size=15, bold=True, color=color)
        add_text(s, detail, x + 0.25, iy + 0.6, 3.5, 0.6, size=12, color=LIGHT_GRAY)

# ── SLIDE 5: Lehman's Laws ────────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "LEHMAN'S LAWS OF SOFTWARE EVOLUTION", 0.5, 0.18, 12, 0.7, size=28, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 6.0, 0.04, GOLD)

laws = [
    ("Law of Continuing Change", "🔄", "CLI → Flask Web App → PyQt5 Desktop → Streamlit. The system continuously evolved to remain useful across different user environments.", RED),
    ("Law of Increasing Complexity", "📈", "Adding features (3D cards, audio, AI, fold mechanic) increased complexity. Managed via OOP refactoring and API separation.", GOLD),
    ("Law of Continuing Growth", "🌱", "Functional content grew each sprint: betting streets, fold mechanics, statistics dashboard, provably-fair audit logs.", GREEN),
    ("Law of Conservation of Familiarity", "⚖️", "Core poker rules never changed. The game engine (HandEvaluator) remained stable while the delivery layer evolved.", PURPLE),
]
for i, (title, icon, desc, color) in enumerate(laws):
    row = i % 2; col = i // 2
    x = 0.5 + col * 6.4; y = 1.3 + row * 2.8
    add_rect(s, x, y, 6.0, 2.6, RGBColor(0x0C,0x10,0x22))
    add_text(s, icon, x + 0.2, y + 0.2, 0.8, 0.8, size=28)
    add_text(s, title, x + 1.1, y + 0.25, 4.7, 0.5, size=16, bold=True, color=color)
    add_text(s, desc, x + 0.2, y + 1.0, 5.6, 1.4, size=13, color=LIGHT_GRAY)

# ── SLIDE 6: Refactoring ─────────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "REFACTORING LEGACY CODE", 0.5, 0.18, 12, 0.7, size=32, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 4.0, 0.04, GOLD)

add_rect(s, 0.4, 1.2, 5.8, 5.8, RGBColor(0x1A,0x06,0x06))
add_text(s, "❌  LEGACY (CLI)", 0.6, 1.3, 5.4, 0.5, size=18, bold=True, color=RED)
legacy = "• Single 1,000+ line script\n• Hardcoded input() loops\n• print() for all UI output\n• Game logic mixed with display\n• No error handling\n• Zero modularity or reuse\n• No database abstraction\n• Untestable spaghetti code"
add_text(s, legacy, 0.6, 1.9, 5.4, 4.8, size=14, color=RGBColor(0xFF,0x99,0x99))

add_rect(s, 7.0, 1.2, 5.8, 5.8, RGBColor(0x03,0x16,0x0C))
add_text(s, "✅  REFACTORED (OOP + API)", 7.2, 1.3, 5.4, 0.5, size=18, bold=True, color=GREEN)
refactored = "• Modular classes: Card, Deck,\n  HandEvaluator, Database, AIPlayer\n• Stateless REST API endpoints\n• JSON responses for all state\n• Context managers for DB safety\n• Try/except exception handling\n• 100% unit-testable engine\n• Multi-platform deployment"
add_text(s, refactored, 7.2, 1.9, 5.4, 4.8, size=14, color=RGBColor(0x86,0xEF,0xAC))

add_text(s, "→", 6.0, 3.6, 1.2, 1.0, size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── SLIDE 7: Testing & CI/CD ─────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "TESTING & AUTOMATED CI/CD", 0.5, 0.18, 12, 0.7, size=32, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 4.5, 0.04, GOLD)

tests = [
    ("🧪 Unit Tests\n(test_gravity_poker.py)", "• test_card_creation()\n  Verifies Card object attributes\n• test_deck_shuffle_and_draw()\n  52 cards, draw reduces count\n• test_hand_evaluator_royal_flush()\n  Asserts Royal Flush = rank 9\n• test_ai_player_decision()\n  AI calls/raises on high score", PURPLE),
    ("⚙️ Automated CI/CD\n(.github/workflows/test.yml)", "• Triggers on every git push\n  and pull_request to main\n• Spins up Ubuntu environment\n• Installs Flask, PyQt5 deps\n• Runs full unittest suite\n• Blocks merge if tests fail\n• Prevents regression bugs", CYAN),
    ("🛡️ Exception Handling", "• SQLite: closing() context managers\n  prevent DB lock leaks\n• Flask: try/except around\n  all DB write operations\n• JS Frontend: try/catch on\n  all fetch() API calls\n• Errors logged to Audit Log\n  instead of crashing UI", GREEN),
]
for i, (title, body, color) in enumerate(tests):
    x = 0.4 + i * 4.3
    add_rect(s, x, 1.2, 4.1, 5.8, RGBColor(0x0C,0x10,0x22))
    add_rect(s, x, 1.2, 4.1, 0.05, color)
    add_text(s, title, x + 0.2, 1.3, 3.7, 0.9, size=15, bold=True, color=color)
    add_text(s, body, x + 0.2, 2.3, 3.7, 4.5, size=13, color=LIGHT_GRAY)

# ── SLIDE 8: Deployment ───────────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "MULTI-PLATFORM DEPLOYMENT", 0.5, 0.18, 12, 0.7, size=32, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 4.0, 0.04, GOLD)

platforms = [
    ("🌐", "Flask Web App", "gravity_poker.py", "py gravity_poker.py", "HTML/CSS/JS frontend. Serves REST API. Runs on any browser. Deployable to Heroku/WSGI.", PURPLE),
    ("🖥️", "PyQt5 Desktop", "desktop_app.py", "py desktop_app.py", "Bundles Chromium engine. Exact web UI replica. Polls Flask server internally. No browser needed.", CYAN),
    ("📊", "Streamlit App", "streamlit_app.py", "py -m streamlit run ...", "Pure Python web UI. No JS/CSS files. Rapid dashboard deployment. Auto-reloads on code save.", GREEN),
]
for i, (icon, name, file, cmd, desc, color) in enumerate(platforms):
    x = 0.4 + i * 4.3
    add_rect(s, x, 1.2, 4.1, 6.0, RGBColor(0x0C,0x10,0x22))
    add_text(s, icon, x + 0.15, 1.3, 0.8, 0.8, size=32)
    add_text(s, name, x + 1.0, 1.45, 2.9, 0.5, size=20, bold=True, color=color)
    add_rect(s, x + 0.2, 2.2, 3.7, 0.04, color)
    add_text(s, "File:", x + 0.2, 2.35, 1.0, 0.35, size=12, color=LIGHT_GRAY)
    add_text(s, file, x + 0.9, 2.35, 2.9, 0.35, size=12, bold=True, color=WHITE)
    add_text(s, "Run:", x + 0.2, 2.75, 1.0, 0.35, size=12, color=LIGHT_GRAY)
    add_text(s, cmd, x + 0.9, 2.75, 3.1, 0.35, size=11, bold=True, color=GOLD)
    add_rect(s, x + 0.2, 3.2, 3.7, 0.04, RGBColor(0x20,0x28,0x48))
    add_text(s, desc, x + 0.2, 3.35, 3.7, 3.6, size=13, color=LIGHT_GRAY)

# ── SLIDE 9: Peer Reviews & SPI ──────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "PEER REVIEWS & SOFTWARE PROCESS IMPROVEMENT", 0.5, 0.18, 12, 0.7, size=26, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 5.0, 0.04, GOLD)

reviews = [
    ("🔍 Code Inspection", "CSS Inheritance Bug", "Red suits (♥ ♦) were rendering black.\nInspection found .card-front had hardcoded color: #1a1b24 overriding .suit-red.\nFixed by removing hardcoded color.", RED),
    ("🔄 Walkthrough", "Desktop Rendering Failure", "pywebview silently failed on Windows (missing Edge WebView2 runtime).\nWalkthrough identified ENV dependency.\nPivoted to PyQt5 + QWebEngineView.", GOLD),
    ("📈 SPI (Deming Cycle)", "Process Improvement", "PLAN: Define REST API contracts.\nDO: Implement modular classes.\nCHECK: Run unit tests on CI.\nACT: Refactor on regression failures.", GREEN),
    ("👥 Team Roles", "Contribution Model", "Scrum Master: Timeline & Lehman's.\nBackend: Flask API + DB schema.\nFrontend: CSS/JS + PyQt5 wrapper.\nQA: test_gravity_poker.py + CI/CD.", PURPLE),
]
for i, (rtype, title, desc, color) in enumerate(reviews):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.4; y = 1.3 + row * 2.9
    add_rect(s, x, y, 6.2, 2.6, RGBColor(0x0C,0x10,0x22))
    add_rect(s, x, y, 6.2, 0.45, color)
    add_text(s, rtype, x + 0.15, y + 0.07, 5.8, 0.35, size=13, bold=True, color=DARK_BG)
    add_text(s, title, x + 0.15, y + 0.58, 5.8, 0.4, size=16, bold=True, color=color)
    add_text(s, desc, x + 0.15, y + 1.05, 5.8, 1.35, size=12, color=LIGHT_GRAY)

# ── SLIDE 10: Key Features Demo ─────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "KEY FEATURES", 0.5, 0.18, 12, 0.7, size=32, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 2.5, 0.04, GOLD)

features = [
    ("🃏", "3D Card Flip Animation", "CSS rotateY(180deg) transform with cubic-bezier easing. Cards deal with staggered 120ms delays."),
    ("🔊", "Web Audio Synthesis", "Zero-dependency sound engine using Web Audio API OscillatorNode. Win/loss/deal jingles."),
    ("✨", "Confetti Celebration", "Canvas-based particle system with 200 physics particles on player win."),
    ("🏦", "Glassmorphism UI", "backdrop-filter: blur(), rgba backgrounds, and glow effects for a premium casino feel."),
    ("🎰", "FOLD Mechanic", "Full API endpoint /api/fold. Instantly ends hand, reveals dealer cards, logs stats."),
    ("📊", "Luck Dashboard", "Real-time win rate meter, streak tracker, segmented bar chart from SQLite stats."),
]
cols = 3
for i, (icon, title, desc) in enumerate(features):
    col = i % cols; row = i // cols
    x = 0.4 + col * 4.3; y = 1.3 + row * 2.8
    add_rect(s, x, y, 4.1, 2.5, RGBColor(0x0C,0x10,0x22))
    add_text(s, icon, x + 0.15, y + 0.2, 0.7, 0.7, size=28)
    add_text(s, title, x + 0.95, y + 0.3, 3.0, 0.5, size=15, bold=True, color=GOLD)
    add_text(s, desc, x + 0.15, y + 1.0, 3.8, 1.35, size=12, color=LIGHT_GRAY)

# ── SLIDE 11: Conclusion ─────────────────────────────────────────────────
s = add_slide()
add_bg(s)
add_rect(s, 0, 0, 13.33, 1.0, RGBColor(0x0F,0x0F,0x28))
add_text(s, "CONCLUSION & LEARNING OUTCOMES", 0.5, 0.18, 12, 0.7, size=28, bold=True, color=CYAN)
add_rect(s, 0.5, 0.95, 5.0, 0.04, GOLD)

outcomes = [
    ("✅", "Agile Process", "Sprints enabled fast iteration and mid-project pivots (pywebview → PyQt5).", GREEN),
    ("✅", "Lehman's Laws", "Each deployment tier proves the laws of Continuing Change and Growth.", CYAN),
    ("✅", "Refactoring", "Legacy 1000-line script transformed into a modular, multi-platform product.", PURPLE),
    ("✅", "Testing & CI", "GitHub Actions ensures no regression; unit tests validate the math engine.", GOLD),
    ("✅", "Exception Handling", "DB context managers + JS try/catch prevent silent failures in production.", RED),
    ("✅", "Peer Reviews", "Two documented reviews led to real bug fixes and architectural improvements.", GREEN),
]
for i, (check, title, desc, color) in enumerate(outcomes):
    row = i % 3; col = i // 3
    x = 0.4 + col * 6.4; y = 1.3 + row * 1.9
    add_rect(s, x, y, 6.2, 1.65, RGBColor(0x0C,0x10,0x22))
    add_text(s, check, x + 0.15, y + 0.35, 0.6, 0.6, size=24, color=color)
    add_text(s, title, x + 0.85, y + 0.15, 5.1, 0.45, size=16, bold=True, color=color)
    add_text(s, desc, x + 0.85, y + 0.65, 5.1, 0.8, size=13, color=LIGHT_GRAY)

add_rect(s, 0.5, 7.0, 12.3, 0.04, PURPLE)
add_text(s, "Gravity Flip Poker — Software Construction & Development Final Project", 0, 7.1, 13.33, 0.3, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save
out = r"c:\Users\Lenovo\OneDrive\Documents\New folder\uml\GravityFlipPoker_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
